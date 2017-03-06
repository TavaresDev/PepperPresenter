from naoqi import ALProxy
from pptx import Presentation
from Tkinter import Tk
from tkFileDialog import askopenfilename

robot_ip = "set robot ip" #!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!
my_notes=[]
#define presentation path
Tk().withdraw()
my_path = askopenfilename()

#create a presentation object from the library pptx and get its number of slides
my_pres = Presentation(my_path)
number_of_slides = len(my_pres.slides)

#crete the proxy for speech
#tts = ALProxy("ALTextToSpeech", robot_ip, 9559)
#use the tablet
#tablet = session.service("ALTabletService")
j=0
i=1
#iterate over the number of slides, making pepper show the image and tell notes
for slide in my_pres.slides:
        #tablet show "slides/slide" + str(i+1) + ".jpg"
        #slide_path = "http://198.18.0.1/" + "path¿?" + "slides/slide" + str(i+1) + ".jpg"
        #tablet.show(slide_path)
        #iterate over the slide notes
        print i
        i = i+1
        if slide.has_notes_slide: #chech if we have notes
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                #those 2 following extractions are due to how the XML is structured
                for paragraph in text_frame.paragraphs:
                        print paragraph.text
